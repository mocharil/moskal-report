from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Cm, Pt
import ast
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# === FUNGSI TEXTBOX GENERIK ===
def add_text(slide, text, left_cm, top_cm, width_cm, height_cm, font_name, font_size, bold=True, color=(0, 0, 0), align='left'):
    textbox = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)

    p.alignment = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }.get(align, PP_ALIGN.LEFT)

    return textbox

# === FUNGSI KOTAK BIRU DYNAMIC PANJANG ===
def add_topic_box(slide, topic_text, left_cm, top_cm, padding_cm=0.7):
    full_text = f"• {topic_text}"
    font_size_pt = 24
    approx_char_width = 0.4  # estimasi lebar per karakter dalam cm
    text_length_cm = len(full_text) * approx_char_width + padding_cm * 2
    height_cm = 1.3

    # Kotak biru muda
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left_cm),
        Cm(top_cm),
        Cm(text_length_cm),
        Cm(height_cm)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 242, 255)
    shape.line.fill.background()

    # Text di atasnya
    add_text(
        slide,
        full_text,
        left_cm,
        top_cm,
        text_length_cm,
        height_cm,
        font_name="Arial",
        font_size=font_size_pt,
        bold=True,
        color=(0, 102, 204)
    )

def date_range_box_cover(slide, text, left_cm, top_cm, height_cm=0.8, icon_path=None):
    # Estimasi lebar berdasarkan panjang teks
    padding_cm = 1.2
    text_length_cm = len(text) * 0.15 + padding_cm * 2
    box_width = text_length_cm

    # Kotak abu (rounded rectangle)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(left_cm),
        Cm(top_cm),
        Cm(box_width),
        Cm(height_cm)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # Abu terang
    shape.line.fill.background()

    # Tambahkan ikon kalender jika ada
    if icon_path:
        slide.shapes.add_picture(
            icon_path,
            Cm(left_cm + 0.3),
            Cm(top_cm + 0.15),
            height=Cm(0.38)
        )

    # Tambahkan teks tanggal
    textbox = slide.shapes.add_textbox(
        Cm(left_cm + 0.8),
        Cm(top_cm),
        Cm(box_width - 1.2),
        Cm(height_cm)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos Display"
    font.size = Pt(12)
    font.bold = True
    font.color.rgb = RGBColor(102, 102, 102)

    return shape
    
def add_date_range_box(slide, text, left_cm, top_cm, height_cm=0.8, icon_path=None):

    # Kotak abu
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(27.87),
        Cm(0.83),
        Cm(4.06),
        Cm(0.59)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # Abu terang
    shape.line.fill.background()

    # Ikon kalender
    if icon_path:
        slide.shapes.add_picture(
            icon_path,
            Cm(28.05),
            Cm(0.91),
            width = Cm(0.34),
            height=Cm(0.34)
        )

        
    # Teks
    textbox = slide.shapes.add_textbox(
        Cm(28.3),
        Cm(0.84),
        Cm(4.2),
        Cm(1)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos"
    font.size = Pt(8)
    font.bold = True
    font.color.rgb = RGBColor(102, 102, 102)

    return shape

def add_slide_template(slide, topic="KADIN", range_date="15 Feb 2025 – 18 Mar 2025", 
                       page_num=2, total_pages=12, icon_path=None,
                       padding_cm=0.2):
  
    # Kotak latar belakang KADIN
    full_text = f"• {topic}"
    approx_char_width = 0.19  # estimasi lebar per karakter dalam cm
    text_length_cm = len(full_text) * approx_char_width + padding_cm * 2
    
    left_box = 1.23
    
    
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,  Cm(left_box), Cm(0.83), Cm(text_length_cm),  Cm(0.68))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 242, 255)
    shape.line.fill.background()
    
    
    # Tulisan KADIN
    txBox = slide.shapes.add_textbox(Cm(1.05), Cm(0.83), Cm(text_length_cm), Cm(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = full_text
    font = run.font
    font.name = "Arial"
    font.size = Pt(10)
    font.bold = True
    font.color.rgb = RGBColor(0, 102, 204)

    # Page x of xx
    #txBox2 = slide.shapes.add_textbox(Cm(1.05),Cm(0.82), Cm(4), height)
    #tf2 = txBox2.text_frame
    #p2 = tf2.paragraphs[0]
    #p2.text = f"Page {page_num} of {total_pages}"
    #p2.font.size = Pt(8)
    #p2.font.name = "Arial"
    #p2.font.color.rgb = RGBColor(102, 102, 102)

    # Garis Horizontal
    left_line = left_box + text_length_cm + 0.25
    width_line = 27.57 - left_line
    
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left_line), Cm(1.1), Cm(width_line), Cm(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(160, 160, 160)
    line.line.fill.background()

    # Kotak tanggal (pakai fungsi)
    add_date_range_box(slide, range_date, left_cm=26.81 ,top_cm=0.83, height_cm=0.9, icon_path=icon_path)
    
    title_box = slide.shapes.add_textbox(0, Cm(18.46), Cm(6), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "© Moskal 2025, All rights reserved."
    font = run.font
    font.name = "Aptos Display"
    font.size = Pt(8)
    font.bold = False
    font.color.rgb = RGBColor(140, 140, 140)
   
    
    return slide

def add_chart_title(slide, text, left, top, width, height, font_size, font_style, font_color, bold, wrap=False, url=None):
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    
    # Jika URL disediakan, tambahkan hyperlink ke teks
    if url:
        # Tambahkan run dengan hyperlink
        hyperlink = p.add_run()
        hyperlink.text = text
        hyperlink.hyperlink.address = url
        
        # Format teks hyperlink
        font = hyperlink.font
        font.name = font_style
        font.size = font_size
        font.bold = bold
        font.color.rgb = RGBColor(67,59,211)
        # Opsional: tambahkan garis bawah untuk menunjukkan ini adalah link
        font.underline = False
    else:
        # Perilaku normal tanpa hyperlink
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = font_style
        font.size = font_size
        font.bold = bold
        font.color.rgb = font_color
    
    # Atur word wrap jika diperlukan
    tf.word_wrap = wrap
    
    return title_box    
    
def line_title(slide, left,top,width,height):
    # Garis biru di bawah judul
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left,top,width,height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 112, 255)
    line.line.fill.background()